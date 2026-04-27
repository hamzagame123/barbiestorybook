from pathlib import Path
from shutil import copy2

from pptx import Presentation


ROOT = Path(r"C:\Users\HAMZA\BarbieStorybook")
PPTX = ROOT / "Critical_project_context" / "BarbieStorybook_VisualLanguage.pptx"
BACKUP = ROOT / "Critical_project_context" / "BarbieStorybook_VisualLanguage.backup-before-design-system-update.pptx"


UPDATES = {
    1: {
        6: "Visual\nLanguage",
        7: "A beautiful design system for Barbie StoryScenes:\nicon, audio, UI, and 3D asset language for real-time mobile AR.",
        9: "Immersive Technology - Visual Language Assignment",
    },
    2: {
        2: "Colour, Type, Icon & Tone",
        7: "Page background",
        11: "Gradient depth",
        15: "Primary CTA / focus",
        19: "UI frames / dividers",
        23: "Soft labels / glow",
        27: "Content slide BG",
        31: "Fraunces - serif display\nRounded icon shapes - soft silhouette language",
        32: "DM Mono - system labels / compact UI",
        36: "Dreamy. Directed. Polished.\nIcons stay rounded, toy-like, and readable at mobile size.\nAudio cues should feel sparkly, warm, and brief.",
    },
    3: {
        6: "radial-gradient\ntop -> transparent 36%\n\nlinear-gradient\n#17020f -> #0c0008",
        9: "MOBILE META (index.html)",
        10: "· theme-color: #0C0008",
        11: "· apple-mobile-web-app-capable: yes",
        12: "· status-bar-style: black-translucent",
        13: "· viewport: no-scale, viewport-fit=cover",
        16: "FONT + AUDIO TOKENS",
        17: "Fraunces\nopsz 9-144 · wght 300, 400\nSOFT axis: 0 and 100 for sharp vs rounded tone\n\nDM Mono\nwght 400, 500\n\nAudio cues\nshort sparkle, soft click, warm success chime",
    },
    4: {
        1: "UI COMPONENTS",
        2: "In-AR Overlay Language",
        6: "SCANNING   START AR   BOOK",
        8: "GLIMMER\nStart your scene...\nLIBRARY   MIC\nCAPTURE",
        9: "HUD wireframe",
        12: "SCANNING -> SURFACE state chip, top-left",
        15: "START AR - launches Needle Go App Clip session",
        17: "#book-btn",
        18: "BOOK - opens saved StoryScenes pages",
        20: "#glimmer-btn",
        21: "GLIMMER - opens the narrator / story helper panel",
        23: "#library-btn",
        24: "Library launcher - opens dolls, backdrops, and accessories",
        26: "#mic-btn",
        27: "MIC - short voice capture for scene ideas and captions",
        29: "#style-pill",
        30: "Style chip - shows current scene mood or prompt preset",
        32: "#place-state",
        33: "Placement state - disabled until surface is found",
        35: "#scene-save",
        36: "SAVE SCENE - enabled after doll, backdrop, and props are composed",
        39: "CAPTURE - saves the framed StoryScene moment",
        42: "Hidden toast - progress, save, and audio feedback states",
    },
    5: {
        2: "Scrapbook Visual Language",
        6: "📸  polished\nStoryScene",
        7: "\"She built a little dream world and stepped right in.\"",
        8: "Glimmer caption · Fraunces italic",
        13: "Warm paper + blush tint for keepsake feeling",
        17: "2-column CSS Grid, single-column on narrow mobile",
        21: "Hot pink strip, ~55% opacity, centred on top edge",
        25: "Fraunces italic - editorial serif, whimsical",
        29: "\"Your first StoryScene starts here ✨\" - always visible",
        33: "× glyph, top-right, always accessible",
    },
    6: {
        2: "Real-Time Rendering on Mobile",
        8: "Style\nGuide",
        9: "Glimmer + art direction",
        10: "Colour mood, prop family, silhouette target",
        16: "Reference\nImage",
        17: "Nano Banana",
        18: "Barbie-tone visual reference for doll or backdrop",
        24: "3D\nGeneration",
        25: "Hyper3D Rodin",
        26: "Consistent GLB output for the selected doll concept",
        32: "Asset\nOptimization",
        33: "GLB + textures",
        34: "Mobile-safe geometry, embedded PBR, reduced texture cost",
        40: "Scene\nPlacement",
        41: "Three.js / WebXR",
        42: "Placed in AR with authored lighting, props, and UI",
        45: "MOBILE CONSTRAINT: real-time rendering on iPhone means keeping geometry readable but efficient, textures compact, and draw calls low. Asset quality must stay beautiful without dropping responsiveness, so StoryScenes uses controlled silhouettes, optimized GLB delivery, and consistent material rules.",
    },
    7: {
        6: "main.ts / style.css",
        7: "Deep plum background and cream content shell create strong contrast for Barbie pink UI while keeping AR preview edges clean.",
        11: "main.ts",
        12: "Warm blush directional light defines the hero doll shape and keeps generated characters flattering on mobile.",
        16: "main.ts",
        17: "Soft fill light lifts shadows so props and accessories stay readable without overloading the scene.",
        20: "Reticle + Icons",
        21: "ARPlacement.ts / LibraryUI",
        22: "Rounded pink reticle, soft chips, and bucket icons all use the same toy-like silhouette language for quick mobile recognition.",
        25: "Audio Cues",
        26: "Narrator / capture flow",
        27: "Short sparkle, capture click, and save chime reinforce state changes without cluttering the screen. Audio matches the same polished Barbie tone as the UI.",
        30: "3D Asset Consistency",
        31: "BackgroundSpawner / CharacterSpawner",
        32: "Dolls, props, and backdrops all follow the same polish rules: clean silhouettes, readable colour blocks, soft highlights, and mobile-safe scale.",
    },
    8: {
        2: "What This Assignment Demonstrates",
        8: "Professional Production Pipeline",
        9: "Identify and Create\nAR Components",
        10: "The deck maps a professional production pipeline across UI, icon, audio, and 3D systems: design direction, asset creation, optimization, and live WebXR placement. Each component is designed for real-time mobile delivery, not static mockups.",
        15: "High-Quality Immersive Content",
        16: "3D Models, UI &\nImmersive Capture",
        17: "StoryScenes balances beauty with hardware limits: polished 3D models, branded lighting, readable icons, warm audio feedback, and scrapbook capture that all feel consistent on mobile AR hardware.",
    },
    9: {
        9: "✓  Brand system: colour, type, icon, and audio tone",
        10: "✓  UI component inventory for the immersive experience",
        11: "✓  3D asset language + real-time rendering constraints",
        12: "✓  Asset pipeline for dolls, props, backdrops, and capture",
        13: "✓  Mobile optimization notes for target hardware",
        14: "✓  Learning outcomes tied to production workflow",
        17: "🔗  UNITY PROJECT FOLDER (WEBLINK)",
        18: "Assets / UI / Audio",
        19: "Design system tokens, icons, and cue references",
        20: "Assets / Models / Materials",
        21: "3D dolls, props, backdrops, optimized materials",
        22: "Assets / Scripts / AR",
        23: "Placement, interaction, scene composition",
        24: "Project Settings / Scene setup",
        25: "Lighting, camera, platform configuration",
        26: "Build / Documentation",
        27: "Target hardware notes and export-ready project files",
    },
    10: {
        3: "Every icon, every panel,",
        4: "every prop - one language.",
        5: "Barbie StoryScenes  ·  Building the Visual Language",
    },
}


def update_shape_text(shape, new_text: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    shape.text = new_text


def main() -> None:
    if not BACKUP.exists():
        copy2(PPTX, BACKUP)

    prs = Presentation(PPTX)
    for slide_no, shape_updates in UPDATES.items():
        slide = prs.slides[slide_no - 1]
        for shape_idx, new_text in shape_updates.items():
            update_shape_text(slide.shapes[shape_idx], new_text)

    prs.save(PPTX)
    print(f"Updated: {PPTX}")
    print(f"Backup: {BACKUP}")


if __name__ == "__main__":
    main()
