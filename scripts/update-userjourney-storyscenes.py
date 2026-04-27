from pptx import Presentation


TARGET = r"C:\Users\HAMZA\BarbieStorybook\Critical_project_context\BarbieStorybook_UserJourney_v2.pptx"


REPLACEMENTS = {
    1: {
        8: "How Barbie StoryScenes guides kids through a living AR\nscene-making experience using triggers, not instructions.",
    },
    2: {
        3: "Immersive Design Teaches the Scene",
        7: "A trigger is any visual, spatial, or auditory cue that tells the child what to do next without a manual, tooltip, or instruction screen. In StoryScenes, the environment, the HUD, and Glimmer all work together to teach the flow.",
        12: "Reticle appears on a real surface — shows exactly where the StoryScene can be placed",
        17: "Surface badge switches ● SCANNING → ✓ SURFACE — confirms the room is ready for scene placement",
        22: "OPEN LIBRARY becomes the next clear action after placement — staging happens in a natural sequence",
    },
    3: {
        10: "Open URL in iPhone Safari",
        11: "😮 Curious",
        16: "02 SPARK",
        17: "Tell Glimmer a scene idea",
        18: "✨ Inspired",
        23: "03 STAGE",
        24: "Pick a doll, backdrop, or prop from the toybox",
        25: "🧸 Choosing",
        30: "04 PLACE",
        31: "Tap a surface to anchor the StoryScene in AR",
        32: "🌍 Grounded",
        37: "05 DIRECT",
        38: "Drag, pinch, and rotate to compose the scene",
        39: "🎯 In Control",
        44: "06 CAPTURE",
        45: "Tap CAPTURE — Nano Banana polishes the moment",
        46: "📸 Directed",
        51: "07 KEEP",
        52: "Glimmer captions it, saved to the book",
        53: "💖 Attached",
    },
    4: {
        8: "Detects horizontal surfaces in real space",
        9: "ARPlacement.ts",
        15: "Visual affordance — shows where the StoryScene will land",
        16: "ARPlacement.ts",
        22: "SCANNING → SURFACE → PLACED state feedback in HUD",
        23: "HUD.ts",
        29: "Shared anchor for toy + stage; drag/pinch to compose",
        30: "SceneRig.ts + SceneGestures.ts",
        35: "Library UI",
        36: "Opens dolls, backdrops, and props for fast scene staging",
        37: "LibraryUI.ts",
        42: "Character Spawner",
        43: "Loads preset or generated Barbie GLBs into the toy slot",
        44: "CharacterSpawner.ts",
        49: "Background Spawner",
        50: "Loads pano backdrops / environment spheres for StoryScenes worlds",
        51: "BackgroundSpawner.ts",
        56: "Book Capture + Storage",
        57: "Screenshots the scene, polishes it, captions it, and stores it locally",
        58: "HUD.ts + ScrapbookStore.ts",
    },
    5: {
        5: "⚠  User doesn't know what kind of scene to make",
        9: "Glimmer gives a spark prompt and short scene beats — the child never starts from a blank page",
        11: "⚠  Doll or backdrop generation takes a few seconds — user thinks app froze",
        15: "Progressive toast messages like 'Loading toy…' and 'Painting backdrop…' communicate that the StoryScene is still being built",
        17: "⚠  Too many actions appear at once",
        21: "The flow is gated in sequence — scan, place, open library, direct, capture — so the child only sees the next meaningful action",
        23: "⚠  User doesn't know AR is available on iPhone",
        27: "Native Needle WebXR button triggers the App Clip flow — Needle Go launches instantly, no installation required",
    },
    6: {
        10: "NEXT ACTION",
        14: "BOOK",
        20: "No surface detected",
        21: "● SCANNING\namber",
        22: "Disabled\n(awaiting surface)",
        23: "Disabled",
        24: "Always active",
        30: "Surface found",
        31: "✓ SURFACE\ngreen",
        32: "PLACE SCENE",
        33: "Disabled",
        34: "Always active",
        40: "Scene placed",
        41: "PLACED\npink",
        42: "OPEN LIBRARY",
        43: "Disabled",
        44: "Always active",
        50: "Staging scene",
        51: "PLACED\npink",
        52: "PICK TOY /\nBACKDROP",
        53: "Disabled",
        54: "Always active",
        60: "Directing (drag/pinch)",
        61: "PLACED\npink",
        62: "ENABLED",
        63: "CAPTURE\nwhite",
        64: "Always active",
        70: "Capturing",
        71: "PLACED\npink",
        72: "Disabled\n(toast: saving…)",
        73: "Disabled\n(active)",
        74: "Always active",
    },
    7: {
        10: "StoryScenes System Design",
        11: "Specific AR + scene systems\nfor the user flow",
        12: "ARPlacement, HUD, SceneRig, SceneGestures, LibraryUI, CharacterSpawner, BackgroundSpawner, NarratorClient, and ScrapbookStore are each mapped to the moment in the journey where they guide or support the child.",
        17: "UX / UI Basics",
        18: "User intentions through\nprofessional UX principles",
        19: "Progressive disclosure, state feedback, spatial affordance, story prompting through Glimmer, and direct manipulation all work together so the child feels like a scene director instead of a user filling out a form.",
    },
    8: {
        10: "✓  This deck — updated for Barbie StoryScenes",
        12: "✓  End-to-end 7-step StoryScenes journey map",
        13: "✓  AR + scene system map from the current prototype",
        14: "✓  4 friction points with StoryScenes resolutions",
        20: "Trigger taxonomy diagram",
        22: "Full StoryScenes journey map",
        24: "UI state map — StoryScenes flow",
        26: "Annotated StoryScenes HUD wireframe",
        28: "AR + scene system dependency chart",
    },
    9: {
        4: "The best interface",
        5: "is the one that",
        6: "invites play.",
        7: "Barbie StoryScenes  ·  Mapping the User Journey",
    },
}


def main():
    prs = Presentation(TARGET)
    for slide_index, shape_map in REPLACEMENTS.items():
        slide = prs.slides[slide_index - 1]
        for shape_index, new_text in shape_map.items():
            slide.shapes[shape_index - 1].text = new_text
    prs.save(TARGET)


if __name__ == "__main__":
    main()
